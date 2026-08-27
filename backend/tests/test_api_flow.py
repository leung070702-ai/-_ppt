from __future__ import annotations

import io
import time

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

from app.main import app


def _pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "项目痛点与价值主张"
    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()


def test_upload_analyze_and_export_without_overwriting_source() -> None:
    source_bytes = _pptx_bytes()
    with TestClient(app) as client:
        response = client.post(
            "/api/projects",
            files={"pptx": ("fixture.pptx", source_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            data={"rubric_text": "商业价值、技术表达、内容结构、视觉呈现"},
        )
        assert response.status_code == 200
        payload = response.json()
        time.sleep(0.2)
        suggestions = client.get(f"/api/jobs/{payload['job_id']}/suggestions").json()
        assert suggestions
        accepted = [suggestions[0]["id"]]
        revision = client.post(f"/api/jobs/{payload['job_id']}/revisions", json={"suggestion_ids": accepted})
        assert revision.status_code == 200
        time.sleep(0.2)
        exported = client.get(f"/api/jobs/{payload['job_id']}/download")
        assert exported.status_code == 200
        original = client.get(f"/api/projects/{payload['project_id']}/download/original")
        assert original.status_code == 200
        assert original.content == source_bytes


def test_rejects_path_traversal_filename() -> None:
    with TestClient(app) as client:
        response = client.post(
            "/api/projects",
            files={"pptx": ("..\\escape.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        assert response.status_code == 400
