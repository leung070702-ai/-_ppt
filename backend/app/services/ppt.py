from __future__ import annotations

import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation


def validate_pptx(path: Path, max_bytes: int = 50 * 1024 * 1024) -> dict[str, Any]:
    result: dict[str, Any] = {"valid": False, "warnings": [], "errors": [], "slide_count": None}
    if path.suffix.lower() != ".pptx":
        result["errors"].append("仅支持 .pptx 文件")
        return result
    if path.stat().st_size > max_bytes:
        result["errors"].append("文件大小不能超过 50 MB")
        return result
    try:
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            if "[Content_Types].xml" not in names or "ppt/presentation.xml" not in names:
                result["errors"].append("文件不是有效的 PPTX 压缩包")
            if any(name.lower().endswith("vbaProject.bin".lower()) for name in names):
                result["warnings"].append("检测到宏项目，导出时不会复制宏代码")
            if any("externalLinks" in name for name in names):
                result["warnings"].append("检测到外部链接，建议导出后人工复核")
        if result["errors"]:
            return result
        presentation = Presentation(str(path))
        result["slide_count"] = len(presentation.slides)
        if result["slide_count"] > 100:
            result["errors"].append("幻灯片数量不能超过 100 页")
            return result
        result["valid"] = True
        return result
    except Exception as exc:  # pragma: no cover - defensive boundary
        result["errors"].append(f"无法读取 PPTX：{exc}")
        return result


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def parse_pptx(path: Path) -> list[dict[str, Any]]:
    presentation = Presentation(str(path))
    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        elements = []
        texts: list[str] = []
        for shape_index, shape in enumerate(slide.shapes):
            text = _shape_text(shape)
            if text:
                texts.append(text)
            elements.append({
                "index": shape_index,
                "type": shape.shape_type.name if hasattr(shape.shape_type, "name") else str(shape.shape_type),
                "text": text[:500],
                "left": round(shape.left / 914400, 2),
                "top": round(shape.top / 914400, 2),
                "width": round(shape.width / 914400, 2),
                "height": round(shape.height / 914400, 2),
            })
        title = texts[0][:80] if texts else f"第 {index} 页"
        slides.append({
            "slide_number": index,
            "title": title,
            "notes": "",
            "elements": elements,
            "text": "\n".join(texts),
        })
    return slides


def apply_actions(source: Path, destination: Path, actions: list[dict[str, Any]]) -> dict[str, Any]:
    prs = Presentation(str(source))
    applied = 0
    skipped = 0
    for action in actions:
        slide_number = action.get("slide_number", 1)
        if not 1 <= slide_number <= len(prs.slides):
            skipped += 1
            continue
        slide = prs.slides[slide_number - 1]
        kind = action.get("kind")
        if kind == "text_replace":
            old = action.get("old", "")
            new = action.get("new", "")
            changed = False
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        if old and old in paragraph.text:
                            for run in paragraph.runs:
                                run.text = run.text.replace(old, new)
                            changed = True
            if changed:
                applied += 1
            else:
                skipped += 1
        elif kind == "font_size":
            size = action.get("size_pt", 24)
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = __import__("pptx").util.Pt(size)
            applied += 1
        else:
            skipped += 1
    destination.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destination))
    return {"applied": applied, "skipped": skipped}
